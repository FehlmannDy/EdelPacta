from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETTE ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0xED, 0xE8, 0xDC)
CARD     = RGBColor(0xFA, 0xF7, 0xF2)
BURGUNDY = RGBColor(0x6B, 0x17, 0x28)
GOLD     = RGBColor(0xB8, 0x96, 0x3E)
BORDER   = RGBColor(0xC8, 0xBF, 0xB2)
TEXT     = RGBColor(0x1A, 0x12, 0x0A)
TEXT2    = RGBColor(0x8A, 0x7A, 0x68)
TEXT3    = RGBColor(0x6B, 0x5A, 0x44)
GREEN    = RGBColor(0x4A, 0x7A, 0x50)
RED      = RGBColor(0x9B, 0x2A, 0x2A)

LOGO = "/tmp/edelpacta_logo.png"


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def r(slide, l, t, w, h, fill, border=None, bw=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def t(slide, text, l, top, w, h, size=16, bold=False, color=TEXT,
      align=PP_ALIGN.LEFT, italic=False, spacing=0):
    box = slide.shapes.add_textbox(Inches(l), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if spacing:
        run._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    return box


def div(slide, l, top, w, color=BORDER, h=0.02):
    r(slide, l, top, w, h, color)


def logo_header(slide, subtitle_label):
    """Consistent header bar across all slides."""
    r(slide, 0, 0, 0.2, 7.5, BURGUNDY)
    r(slide, 0.2, 0, 13.13, 1.15, CARD, border=BORDER, bw=0.5)
    div(slide, 0.2, 1.13, 13.13, BURGUNDY, 0.05)
    slide.shapes.add_picture(LOGO, Inches(0.42), Inches(0.09), Inches(0.88), Inches(0.97))
    t(slide, "EDELPACTA", 1.45, 0.2, 5.0, 0.48,
      size=10, bold=True, color=BURGUNDY, spacing=0.35)
    t(slide, subtitle_label, 1.45, 0.7, 9.0, 0.35,
      size=9, bold=True, color=TEXT2, spacing=0.18)


# ─── PRESENTATION ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
r(s, 0, 0, 0.2, 7.5, BURGUNDY)

# Header
r(s, 0.2, 0, 13.13, 1.15, CARD, border=BORDER, bw=0.5)
div(s, 0.2, 1.13, 13.13, BURGUNDY, 0.05)
s.shapes.add_picture(LOGO, Inches(0.42), Inches(0.09), Inches(0.88), Inches(0.97))
t(s, "EDELPACTA", 1.45, 0.2, 5.0, 0.48, size=10, bold=True, color=BURGUNDY, spacing=0.35)

# Hackathon badge
r(s, 9.6, 0.18, 3.5, 0.78, BG, border=GOLD, bw=0.75)
t(s, "XRPL HACKATHON 2025", 9.65, 0.24, 3.4, 0.34,
  size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER, spacing=0.12)
t(s, "Programmability Track", 9.65, 0.6, 3.4, 0.3,
  size=10, color=TEXT3, align=PP_ALIGN.CENTER, italic=True)

# Main title
t(s, "Real Estate Tokenization\n& Smart Escrow", 0.55, 1.3, 9.5, 2.8,
  size=52, bold=True, color=TEXT)

# Tagline
div(s, 0.55, 4.25, 7.5, GOLD, 0.03)
t(s, "Pacta sunt servanda. By code.",
  0.55, 4.42, 9.5, 0.75, size=22, color=TEXT3, italic=True)

# Creators
r(s, 0.55, 5.38, 7.0, 1.05, CARD, border=BORDER, bw=0.75)
r(s, 0.55, 5.38, 0.08, 1.05, BURGUNDY)
t(s, "CREATED BY", 0.82, 5.5, 3.5, 0.3,
  size=8, bold=True, color=TEXT2, spacing=0.18)
t(s, "Dylan Fehlmann  &  Hugo Germano",
  0.82, 5.82, 6.5, 0.52, size=20, bold=True, color=TEXT)

# Right — 3 stat cards
for i, (num, label) in enumerate([("< 0.01 CHF","Transaction\nCost"),("3 sec","Settlement\nTime"),("0","Counterparty\nRisk")]):
    x = 8.35 + i * 1.65
    r(s, x, 1.32, 1.48, 2.55, CARD, border=BORDER, bw=0.75)
    r(s, x, 1.32, 1.48, 0.05, BURGUNDY)
    t(s, num,   x, 1.55, 1.48, 0.85, size=18, bold=True,
      color=BURGUNDY, align=PP_ALIGN.CENTER)
    t(s, label, x, 2.45, 1.48, 0.85, size=11, color=TEXT2,
      align=PP_ALIGN.CENTER)

div(s, 0.2, 7.2, 13.13, BORDER, 0.02)
t(s, "XRPL  |  WASM Hooks  |  XLS-20 NFT  |  Swiyu e-ID (Swiss Gov)",
  0.5, 7.26, 12.5, 0.28, size=9, color=TEXT2, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE USER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "THE USER")

t(s, "Who suffers from\nthis broken process?",
  0.55, 1.28, 10.0, 1.72, size=40, bold=True, color=TEXT)
div(s, 0.55, 3.05, 11.5, GOLD, 0.03)

# Market stat — full-width banner
r(s, 0.45, 3.25, 12.0, 0.88, CARD, border=BORDER, bw=0.75)
r(s, 0.45, 3.25, 0.08, 0.88, GOLD)
t(s, "Switzerland processes over", 0.72, 3.36, 4.5, 0.36, size=14, color=TEXT2, italic=True)
t(s, "CHF 70 billion", 5.0, 3.32, 3.5, 0.46, size=22, bold=True, color=BURGUNDY)
t(s, "in real estate transactions every year", 8.35, 3.36, 3.9, 0.36, size=14, color=TEXT2, italic=True)

# 3 user personas
personas = [
    (BURGUNDY, "The Buyer",
     "Lock up millions of CHF\nfor 8 weeks in a\nnotary's bank account.",
     "No visibility. Blind trust."),
    (GOLD,     "The Seller",
     "Submit paper documents,\nprove tax compliance,\nwait for manual approval.",
     "Slow. Error-prone. Costly."),
    (BURGUNDY, "The Notary",
     "Manage bank escrows,\nchase documents,\nbear legal liability.",
     "Manual. Outdated. Risky."),
]
pw = 3.75
for i, (col, role, pain, verdict) in enumerate(personas):
    x = 0.45 + i * (pw + 0.37)
    r(s, x, 4.3, pw, 2.88, CARD, border=BORDER, bw=0.75)
    r(s, x, 4.3, pw, 0.06, col)
    t(s, role,    x+0.22, 4.44, pw-0.4, 0.52, size=20, bold=True, color=col)
    div(s, x+0.22, 5.02, pw-0.44, BORDER)
    t(s, pain,    x+0.22, 5.16, pw-0.4, 1.18, size=14, color=TEXT3)
    t(s, verdict, x+0.22, 6.38, pw-0.4, 0.42, size=13, bold=True, color=TEXT2, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "THE PROBLEM")

t(s, "Buying real estate in Switzerland is archaic",
  0.55, 1.3, 11.5, 1.05, size=32, bold=True, color=TEXT)
div(s, 0.55, 2.4, 11.5, GOLD, 0.03)

cards = [
    (BURGUNDY, "2M CHF Locked",
     "Massive capital locked\nin a notary's bank account\nfor up to 8 weeks."),
    (GOLD,     "1–2% in Fees",
     "Up to CHF 40,000\nin notary and banking fees\nper transaction."),
    (BURGUNDY, "Paper-Based KYC",
     "Identity checks still rely\non physical documents\nand manual verification."),
]
cw = 3.75
for i, (col, title, body) in enumerate(cards):
    x = 0.45 + i * (cw + 0.37)
    r(s, x, 2.62, cw, 4.52, CARD, border=BORDER, bw=0.75)
    r(s, x, 2.62, cw, 0.06, col)
    t(s, f"0{i+1}", x+0.22, 2.76, 0.7, 0.92,
      size=40, bold=True, color=col)
    t(s, title, x+0.22, 3.74, cw-0.4, 0.62,
      size=21, bold=True, color=TEXT)
    div(s, x+0.22, 4.44, cw-0.44, BORDER)
    t(s, body, x+0.22, 4.6, cw-0.4, 2.2, size=16, color=TEXT3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "HOW IT WORKS — ATOMIC SETTLEMENT")

t(s, "Four steps. One atomic transaction.",
  0.55, 1.3, 10.0, 0.85, size=32, bold=True, color=TEXT)
div(s, 0.55, 2.22, 11.5, GOLD, 0.03)

steps = [
    (BURGUNDY, "01", "The KYC",
     "The Seller proves compliance\nvia Swiyu App (eIDAS 2.0).\nOracle generates crypto proof."),
    (GOLD,     "02", "The Asset",
     "The Notary mints the property\nas an XLS-20 NFT on XRPL.\nMetadata stored on IPFS."),
    (BURGUNDY, "03", "The Escrow",
     "The Buyer locks 2M CHF\n(stablecoin) in a\nWASM Smart Escrow."),
    (GOLD,     "04", "The Execution",
     "Notary signs final approval.\nWASM Hook verifies signatures.\nNFT and funds swap atomically."),
]
sw = 2.88
for i, (col, num, title, body) in enumerate(steps):
    x = 0.38 + i * (sw + 0.36)
    r(s, x, 2.42, sw, 4.72, CARD, border=BORDER, bw=0.75)
    r(s, x, 2.42, sw, 0.06, col)
    t(s, num, x+0.2, 2.56, sw-0.3, 0.88,
      size=42, bold=True, color=col)
    t(s, title, x+0.2, 3.5, sw-0.3, 0.6,
      size=19, bold=True, color=TEXT)
    div(s, x+0.2, 4.17, sw-0.4, BORDER)
    t(s, body, x+0.2, 4.32, sw-0.3, 2.58,
      size=15, color=TEXT3)
    if i < 3:
        t(s, "›", x+sw+0.08, 5.0, 0.28, 0.5,
          size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# WASM note at bottom
div(s, 0.38, 7.2, 12.1, BORDER)
t(s, "If any condition fails, the Hook triggers a rollback(). Zero counterparty risk.",
  0.55, 7.26, 12.0, 0.28, size=10, color=TEXT2, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "TECH STACK")

t(s, "Built on open standards,\npowered by XRPL",
  0.55, 1.3, 9.5, 1.75, size=40, bold=True, color=TEXT)
div(s, 0.55, 3.1, 11.5, GOLD, 0.03)

stack = [
    (BURGUNDY, "Blockchain",      "XRP Ledger\nwasm.devnet.rippletest.net"),
    (GOLD,     "Smart Contracts", "WASM Hooks in Rust\nxrpl-wasm-stdlib"),
    (BURGUNDY, "Identity",        "Swiyu — Swiss Gov e-ID\nSD-JWT / eIDAS 2.0"),
    (GOLD,     "Oracle",          "TypeScript backend\nDual cryptographic signature"),
    (BURGUNDY, "Assets",          "XLS-20 NFTs\nIPFS metadata (Kubo)"),
    (GOLD,     "Frontend",        "React.js + xrpl.js\nOtsu Wallet (zero-trust)"),
]
tw = 3.9
th = 2.12
for i, (col, title, body) in enumerate(stack):
    col_idx = i % 3
    row_idx = i // 3
    x = 0.38 + col_idx * (tw + 0.28)
    y = 3.28 + row_idx * (th + 0.22)
    r(s, x, y, tw, th, CARD, border=BORDER, bw=0.75)
    r(s, x, y, 0.08, th, col)
    t(s, title, x+0.22, y+0.18, tw-0.3, 0.58,
      size=18, bold=True, color=col)
    div(s, x+0.22, y+0.82, tw-0.3, BORDER)
    t(s, body, x+0.22, y+0.98, tw-0.3, 0.98,
      size=15, color=TEXT3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — IMPACT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "IMPACT")

t(s, "The numbers speak for themselves",
  0.55, 1.3, 11.0, 0.85, size=36, bold=True, color=TEXT)
div(s, 0.55, 2.22, 11.5, GOLD, 0.03)

impacts = [
    (BURGUNDY, "COST",              "CHF 20,000",    "< CHF 0.01",   "Per transaction"),
    (GOLD,     "SPEED",             "8 weeks",       "3 seconds",    "Settlement time"),
    (GREEN,    "COUNTERPARTY RISK", "Central failure","Zero",        "Cryptographic guarantee"),
]
iw = 3.75
for i, (col, label, before, after, sub) in enumerate(impacts):
    x = 0.45 + i * (iw + 0.37)

    # Full card
    r(s, x, 2.42, iw, 4.72, CARD, border=BORDER, bw=0.75)
    r(s, x, 2.42, iw, 0.06, col)

    # Category label
    t(s, label, x+0.2, 2.56, iw-0.3, 0.38,
      size=10, bold=True, color=col, spacing=0.18)

    # "Before" zone — shaded strip
    r(s, x+0.15, 3.06, iw-0.3, 1.18, BG)
    t(s, "BEFORE", x+0.25, 3.1, iw-0.5, 0.3,
      size=8, bold=True, color=TEXT2, spacing=0.15)
    t(s, before, x+0.25, 3.38, iw-0.5, 0.72,
      size=22, bold=False, color=TEXT2)

    # Arrow
    t(s, "↓", x, 4.3, iw, 0.55,
      size=30, bold=True, color=col, align=PP_ALIGN.CENTER)

    # "After" zone — highlighted
    r(s, x+0.15, 4.88, iw-0.3, 1.45, col)
    t(s, "EDELPACTA", x+0.25, 4.93, iw-0.5, 0.3,
      size=8, bold=True, color=CARD, spacing=0.15)
    t(s, after, x+0.25, 5.22, iw-0.5, 0.98,
      size=30, bold=True, color=CARD)

    # Sub-label
    div(s, x+0.2, 6.44, iw-0.4, BORDER)
    t(s, sub, x+0.2, 6.57, iw-0.3, 0.38,
      size=11, color=TEXT2, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GO-TO-MARKET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
logo_header(s, "GO-TO-MARKET")

t(s, "How we reach them",
  0.55, 1.3, 10.0, 0.85, size=44, bold=True, color=TEXT)
div(s, 0.55, 2.22, 11.5, GOLD, 0.03)

channels = [
    (BURGUNDY,
     "Swiss Notary Associations",
     "Direct B2B",
     "Notaries are the legal gatekeepers of every property transfer in Switzerland.\n"
     "We partner with the Schweizerische Notarenkonferenz to pilot EdelPacta\n"
     "as the official digital escrow standard."),
    (GOLD,
     "Swiyu — Swiss Gov e-ID",
     "Built-in distribution",
     "Our KYC is powered by the Swiss government's own digital identity app.\n"
     "Every Swiss citizen with a Swiyu account is a potential user.\n"
     "Government backing accelerates institutional trust."),
    (GREEN,
     "Swiss Banks & Proptech",
     "B2B2C",
     "Real estate platforms (Homegate, ImmoScout24) and cantonal banks\n"
     "integrate EdelPacta as their settlement layer, reaching buyers\n"
     "and sellers at the point of transaction."),
]
cw = 3.75
for i, (col, channel, model, desc) in enumerate(channels):
    x = 0.45 + i * (cw + 0.37)
    r(s, x, 2.48, cw, 4.72, CARD, border=BORDER, bw=0.75)
    r(s, x, 2.48, cw, 0.06, col)
    # Model badge
    r(s, x+0.2, 2.62, cw-0.4, 0.38, BG)
    t(s, model, x+0.22, 2.66, cw-0.44, 0.3,
      size=10, bold=True, color=col, spacing=0.12)
    t(s, channel, x+0.22, 3.08, cw-0.4, 0.72,
      size=18, bold=True, color=TEXT)
    div(s, x+0.22, 3.86, cw-0.44, BORDER)
    t(s, desc, x+0.22, 4.0, cw-0.4, 3.0,
      size=13, color=TEXT3)

# Footer
r(s, 0.2, 7.1, 13.13, 0.4, BURGUNDY)
t(s, "Dylan Fehlmann  &  Hugo Germano  |  XRPL Hackathon 2025  |  Programmability Track",
  0.2, 7.14, 13.13, 0.32, size=11, bold=True, color=CARD, align=PP_ALIGN.CENTER)


# ─── SAVE ─────────────────────────────────────────────────────────────────────
out = "/home/daydozkosmos/Hackathon/XRPL/EdelPacta/EdelPacta_Hackathon_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
